"""
Extraction de texte pour tous les formats autorisés en pièces jointes FluxMin.
PDF / images → OCR (avec fallback texte natif PDF).
Office / OpenDocument / texte → extraction native.
"""
from __future__ import annotations

import io
import re
import zipfile
from typing import Any, Dict, List, Tuple
from xml.etree import ElementTree as ET

# Aligné sur backend/src/common/files/storage.util.ts ALLOWED_UPLOAD_EXTENSIONS
SUPPORTED_EXTENSIONS = frozenset(
    {
        ".pdf",
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".webp",
        ".txt",
        ".csv",
        ".rtf",
        ".doc",
        ".docx",
        ".xls",
        ".xlsx",
        ".ppt",
        ".pptx",
        ".odt",
        ".ods",
        ".odp",
    }
)

IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp"})
TEXT_EXTS = frozenset({".txt", ".csv"})

def _ext(filename: str) -> str:
    name = (filename or "").lower()
    if "." not in name:
        return ""
    return "." + name.rsplit(".", 1)[-1]

def _pages(text: str) -> Dict[str, Any]:
    cleaned = (text or "").strip()
    return {
        "pages": [{"page": 1, "text": cleaned}],
        "texteBrut": cleaned,
    }

def _decode_bytes(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")

def _strip_rtf(data: bytes) -> str:
    raw = _decode_bytes(data)
    try:
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(raw)
    except Exception:
        # Fallback grossier : retirer les groupes RTF
        text = re.sub(r"\\[a-z]+\d* ?", " ", raw)
        text = re.sub(r"[{}]", " ", text)
        return re.sub(r"\s+", " ", text).strip()

def _extract_pdf(data: bytes, filename: str) -> Tuple[str, float, str]:
    """Texte natif PDF si possible, sinon OCR."""
    native = ""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        parts: List[str] = []
        for i, page in enumerate(reader.pages):
            t = page.extract_text() or ""
            if t.strip():
                parts.append(f"--- PAGE {i + 1} ---\n{t}")
        native = "\n".join(parts).strip()
    except Exception:
        native = ""

    if len(native) >= 40:
        return native, 90.0, "pdf_text"

    # PDF scanné / vide → OCR
    from app.services.ocr_service import perform_ocr

    ocr = perform_ocr(data, filename if filename.lower().endswith(".pdf") else "document.pdf")
    return ocr["texteExtrait"]["texteBrut"], float(ocr.get("scoreConfiance") or 80.0), "pdf_ocr"

def _extract_image(data: bytes, filename: str) -> Tuple[str, float, str]:
    from app.services.ocr_service import perform_ocr

    ocr = perform_ocr(data, filename)
    return ocr["texteExtrait"]["texteBrut"], float(ocr.get("scoreConfiance") or 85.0), "ocr"

def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)

def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: List[str] = []
    for sheet in wb.worksheets:
        parts.append(f"--- Feuille : {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)

def _extract_xls(data: bytes) -> str:
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    parts: List[str] = []
    for sheet in book.sheets():
        parts.append(f"--- Feuille : {sheet.name} ---")
        for r in range(sheet.nrows):
            cells = [str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols)]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)

def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"--- Diapositive {i} ---\n" + "\n".join(slide_texts))
    return "\n".join(parts)

def _odf_text_from_zip(data: bytes, content_xml: str = "content.xml") -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        xml_bytes = zf.read(content_xml)
    root = ET.fromstring(xml_bytes)
    texts: List[str] = []
    for node in root.iter():
        if node.text and node.text.strip():
            texts.append(node.text.strip())
        if node.tail and node.tail.strip():
            texts.append(node.tail.strip())
    # Dédupliquer lignes trop répétitives
    seen = set()
    out: List[str] = []
    for t in texts:
        if t not in seen:
            seen.add(t)
            out.append(t)
    return "\n".join(out)

def _extract_ole_printable(data: bytes, min_len: int = 5) -> str:
    """Fallback pour .doc / .ppt binaires : chaînes imprimables."""
    # UTF-16LE words often used in OLE
    try:
        as_utf16 = data.decode("utf-16-le", errors="ignore")
        chunks = re.findall(r"[\wÀ-ÿ][\wÀ-ÿ \t.,;:!?()'\"\-/]{%d,}" % (min_len - 1), as_utf16)
    except Exception:
        chunks = []
    ascii_chunks = re.findall(
        rb"[\x20-\x7E\xC0-\xFF]{%d,}" % min_len,
        data,
    )
    ascii_text = [c.decode("latin-1", errors="ignore") for c in ascii_chunks]
    merged = chunks + ascii_text
    # Filtrer bruit OLE
    filtered = [
        t.strip()
        for t in merged
        if len(t.strip()) >= min_len
        and not t.startswith(("Microsoft", "Root Entry", "WordDocument"))
        and sum(ch.isalpha() for ch in t) >= max(3, len(t) // 4)
    ]
    # Limiter volume
    unique: List[str] = []
    seen = set()
    for t in filtered:
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(t)
        if len(unique) > 400:
            break
    return "\n".join(unique)

def _extract_doc(data: bytes) -> str:
    # Tentative via antiword si disponible
    try:
        import shutil
        import subprocess
        import tempfile
        import os

        if shutil.which("antiword"):
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                tmp.write(data)
                path = tmp.name
            try:
                out = subprocess.check_output(["antiword", path], stderr=subprocess.DEVNULL)
                text = out.decode("utf-8", errors="ignore").strip()
                if text:
                    return text
            finally:
                try:
                    os.unlink(path)
                except OSError:
                    pass
    except Exception:
        pass
    return _extract_ole_printable(data)

def _extract_ppt(data: bytes) -> str:
    return _extract_ole_printable(data)

def extract_document(file_bytes: bytes, filename: str) -> Dict[str, Any]:
    """
    Retourne :
      texteExtrait, langue, scoreConfiance, methodeExtraction, format
    """
    ext = _ext(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Format « {ext or 'inconnu'} » non supporté. "
            f"Formats acceptés : {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    methode = "native"
    score = 88.0
    text = ""

    if ext == ".pdf":
        text, score, methode = _extract_pdf(file_bytes, filename)
    elif ext in IMAGE_EXTS:
        text, score, methode = _extract_image(file_bytes, filename)
    elif ext in TEXT_EXTS:
        text = _decode_bytes(file_bytes)
        score = 95.0
        methode = "text"
    elif ext == ".rtf":
        text = _strip_rtf(file_bytes)
        score = 90.0
        methode = "rtf"
    elif ext == ".docx":
        text = _extract_docx(file_bytes)
        score = 93.0
        methode = "docx"
    elif ext == ".xlsx":
        text = _extract_xlsx(file_bytes)
        score = 93.0
        methode = "xlsx"
    elif ext == ".xls":
        text = _extract_xls(file_bytes)
        score = 88.0
        methode = "xls"
    elif ext == ".pptx":
        text = _extract_pptx(file_bytes)
        score = 92.0
        methode = "pptx"
    elif ext in {".odt", ".ods", ".odp"}:
        text = _odf_text_from_zip(file_bytes)
        score = 90.0
        methode = ext.lstrip(".")
    elif ext == ".doc":
        text = _extract_doc(file_bytes)
        score = 70.0 if text else 40.0
        methode = "doc_legacy"
    elif ext == ".ppt":
        text = _extract_ppt(file_bytes)
        score = 65.0 if text else 40.0
        methode = "ppt_legacy"
    else:
        raise ValueError(f"Aucune stratégie d'extraction pour {ext}")

    text = (text or "").strip()
    if not text:
        raise ValueError(
            f"Aucun texte extractible depuis « {filename} » "
            f"(méthode {methode}). Vérifiez le contenu ou convertissez en PDF/TXT."
        )

    return {
        "texteExtrait": _pages(text),
        "langue": "fr",
        "scoreConfiance": score,
        "methodeExtraction": methode,
        "format": ext,
    }
